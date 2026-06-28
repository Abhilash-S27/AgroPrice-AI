"""
Merge all 4 AgroPrice AI PPTX sections into one final presentation.
Adds fade transitions to every slide.
"""
import copy
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

BASE = r"D:\mca\2nd semester\GEN AI\EL Project\AgroPrice-AI\presentation"

INPUT_FILES = [
    rf"{BASE}\AgroPrice_AI_Slides_1_10.pptx",
    rf"{BASE}\AgroPrice_AI_Slides_11_19.pptx",
    rf"{BASE}\AgroPrice_AI_Slides_20_29.pptx",
    rf"{BASE}\AgroPrice_AI_Slides_30_39.pptx",
]
OUTPUT_FILE = rf"{BASE}\AgroPrice_AI_FINAL.pptx"


def copy_slide(dest_prs, src_slide):
    """
    Append src_slide into dest_prs, preserving all shapes and background.
    Uses the blank layout (last layout in the master) as the carrier.
    """
    # Pick the blank-est layout available
    blank = dest_prs.slide_layouts[-1]
    dest_slide = dest_prs.slides.add_slide(blank)

    # ── Clear the new slide's shape tree ──────────────────────────────────────
    sp_tree = dest_slide.shapes._spTree
    for child in list(sp_tree):
        sp_tree.remove(child)

    # ── Copy all shapes from source ────────────────────────────────────────────
    for child in src_slide.shapes._spTree:
        sp_tree.append(copy.deepcopy(child))

    # ── Copy slide background (set via s.background = {color:…} in PptxGenJS) ─
    src_bg = src_slide._element.find(qn("p:bg"))
    if src_bg is not None:
        dest_bg = dest_slide._element.find(qn("p:bg"))
        if dest_bg is not None:
            dest_slide._element.remove(dest_bg)
        # Insert as second child (after cSld namespace declarations)
        dest_slide._element.insert(2, copy.deepcopy(src_bg))

    return dest_slide


def add_fade_transition(slide):
    """Add a medium-speed fade transition to the slide."""
    # Remove any existing transition
    existing = slide._element.find(qn("p:transition"))
    if existing is not None:
        slide._element.remove(existing)

    # Build transition XML
    trans = etree.SubElement(slide._element, qn("p:transition"))
    trans.set("spd", "med")       # medium speed
    trans.set("advClick", "1")    # advance on click
    trans.set("advTm", "0")
    etree.SubElement(trans, qn("p:fade"))


def main():
    print("Opening base presentation …")
    dest = Presentation(INPUT_FILES[0])

    print(f"  OK Loaded {INPUT_FILES[0]}  ({len(dest.slides)} slides)")

    for src_path in INPUT_FILES[1:]:
        src = Presentation(src_path)
        print(f"  Merging {src_path}  ({len(src.slides)} slides) …")
        for i, src_slide in enumerate(src.slides):
            copy_slide(dest, src_slide)
        print(f"  OK Done ({len(dest.slides)} total so far)")

    print("\nApplying fade transitions to all slides …")
    for slide in dest.slides:
        add_fade_transition(slide)
    print(f"  OK {len(dest.slides)} transitions applied")

    print(f"\nSaving to: {OUTPUT_FILE}")
    dest.save(OUTPUT_FILE)
    print(f"DONE  AgroPrice_AI_FINAL.pptx  ({len(dest.slides)} slides) saved successfully.")


if __name__ == "__main__":
    main()
