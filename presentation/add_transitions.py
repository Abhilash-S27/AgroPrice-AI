"""
Add fade transitions to every slide in AgroPrice_AI_FINAL.pptx
"""
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

PATH = r"D:\mca\2nd semester\GEN AI\EL Project\AgroPrice-AI\presentation\AgroPrice_AI_FINAL.pptx"

prs = Presentation(PATH)
print(f"Loaded {len(prs.slides)} slides")

for i, slide in enumerate(prs.slides, 1):
    existing = slide._element.find(qn("p:transition"))
    if existing is not None:
        slide._element.remove(existing)
    trans = etree.SubElement(slide._element, qn("p:transition"))
    trans.set("spd", "med")
    trans.set("advClick", "1")
    etree.SubElement(trans, qn("p:fade"))

prs.save(PATH)
print(f"DONE: fade transitions applied to all {len(prs.slides)} slides")
